import PyPDF2
from pptx import Presentation


def pdf_has_text(pdf_path, keyword):
    pdf_file_obj = open(pdf_path, 'rb')
    pdf_reader = PyPDF2.PdfFileReader(pdf_file_obj)
    for page in range(pdf_reader.numPages):
        page_obj = pdf_reader.getPage(page)
        if keyword in page_obj.extractText():
            return True
    pdf_file_obj.close()
    return False


def pptx_has_text(pptx_path, keyword):
    pptx_file_obj = open(pptx_path, 'rb')
    prs = Presentation(pptx_file_obj)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                if keyword in shape.text:
                    return True
    return False


def get_docs_list(request, queryset):
    docs = []
    keyword = request.GET.get('keyword')
    search_in = request.GET.get('in', 'file_name')
    if search_in == "file_name" and keyword:
        docs = queryset.filter(file__contains=keyword)
    elif search_in == "content":
        has_text = False
        for document in queryset:
            if document.type == "PDF":
                has_text = pdf_has_text(document.file.path, keyword)
            if document.type == "PPTX":
                has_text = pptx_has_text(document.file.path, keyword)
            if has_text:
                docs.append(document)
    return docs
